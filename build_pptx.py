# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

AR = "Tahoma"
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x2E, 0x6F, 0xB0)
ORANGE = RGBColor(0xC0, 0x5A, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title, lines, note=None):
    s = prs.slides.add_slide(BLANK)
    # title bar
    tb = s.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = title
    r.font.name = AR
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = NAVY
    # body
    bx = s.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.9))
    btf = bx.text_frame
    btf.word_wrap = True
    first = True
    for lvl, txt in lines:
        p = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.RIGHT
        p.level = lvl
        r = p.add_run()
        bullet = "•  " if lvl == 0 else "–  "
        r.text = bullet + txt
        r.font.name = AR
        r.font.size = Pt(16 if lvl == 0 else 14)
        r.font.color.rgb = GREY if lvl else RGBColor(0x22, 0x22, 0x22)
    if note:
        nb = s.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.5))
        np = nb.text_frame.paragraphs[0]
        np.alignment = PP_ALIGN.RIGHT
        rr = np.add_run()
        rr.text = note
        rr.font.name = AR
        rr.font.size = Pt(11)
        rr.font.italic = True
        rr.font.color.rgb = ORANGE
    return s


# ---------- Slides ----------
add_slide("نظام تتبّع المواقع والمسير — توثيق شامل", [
    (0, "المشروع: Yarmouk Water Management Pro (إدارة مياه عجلون/يarmouk)."),
    (0, "الأطراف الثلاثة: خادم خلفي (Backend) + تطبيق Flutter على هواتف (المدير، المفوض، الموظف المُتتبَّع)."),
    (0, "كل الاتصالات عبر REST (HTTP/HTTPS) + قناة WebSocket للتنبيهات الفورية."),
    (0, "هذا الملف يشرح المعمارية، نموذج البيانات، نقاط النهاية، تدفّق البيانات، والشبكة بدقة."),
], note="وثيقة نصية فقط — لا تحتوي صوراً.")

add_slide("1) نظرة عامة معمارية", [
    (0, "الخادم: Python + FastAPI + SQLAlchemy (async) + SQLite، يُشغَّل بـ uvicorn على المنفذ 8001."),
    (0, "التطبيق: Flutter (Dart) عربي RTL؛ خرائط flutter_map (بلاطات OpenStreetMap)؛ تصدير PDF عبر مكتبة printing."),
    (0, "المصادقة: JWT Bearer؛ الدخول POST /auth/login بـ {employee_number, password}."),
    (0, "الأدوار: general_manager / hr_manager (يبدأون/يوقفون ويعيّنون المفوض) — موظف عادي (يُتتبَّع) — مفوض (يراقب ويطبع فقط)."),
    (0, "الموظف المفوض: موظف واحد يُحدَّد عالمياً عبر مفتاح AppSetting.tracking_viewer_id."),
])

add_slide("2) الخادم (Backend) — التفاصيل", [
    (0, "المسار: backend/app/main.py — التطبيق app.main:app."),
    (0, "التشغيل: uvicorn app.main:app --host 0.0.0.0 --port 8001 (يُفضّل كخدمة ويندوز دائمة)."),
    (0, "الوحدات: routers/gps.py (التتبع)، services/notifications.py (التنبيهات)، realtime/ws.py (WebSocket)."),
    (0, "القاعدة: SQLite (backend/yarmouk_water_pro.db) تُنشأ وتُبدأ بالبذور (seed) عند الإقلاع."),
    (0, "الوسيط: CORS + RateLimit + تدقيق (audit) لكل الطلبات."),
])

add_slide("3) نموذج البيانات (الجداول)", [
    (0, "employees: بيانات الموظف + geofence_lat / geofence_lng / geofence_radius_m / geofence_exempt (منطقة العمل)."),
    (0, "field_tracking_sessions: جلسة تتبّع — employee_id, started_by_id, track_mode, track_interval, status."),
    (0, "field_tracking_sessions (تكملة): is_outside, outside_started_at, outside_distance_m, last_lat, last_lng, last_point_at."),
    (0, "field_tracking_points: نقطة موقع — session_id, employee_id, latitude, longitude, accuracy, recorded_at."),
    (0, "geofence_breaches: خرق منطقة — breach_id, session_id, employee_id, started_at, ended_at, duration_seconds, distance_m."),
    (0, "app_settings: مفتاح tracking_viewer_id (الموظف المفوض).  notifications: سجل الإشعارات."),
])

add_slide("4) نقاط نهاية GPS (Endpoints)", [
    (0, "POST /gps/start — المدير/الموارد — إنشاء جلسة (target, mode, interval)."),
    (0, "POST /gps/stop — المدير/الموارد — إنهاء الجلسة."),
    (0, "POST /gps/point — الموظف المُتتبَّع — إرسال نقطة + كشف الخروج."),
    (0, "GET /gps/my-active — الموظف — حالة تتبّعه (mode/interval/is_outside/المسافة)."),
    (0, "GET /gps/view — المدير/المفوض — كل الجلسات النشطة + آخر نقطة + geofence + is_outside."),
    (0, "GET /gps/history?employee_id= — المدير/المفوض — نقاط المسير (route)."),
    (0, "GET /gps/employees — المدير/المفوض — دليل الموظفين مع حقول geofence."),
    (0, "POST /gps/set-viewer + GET /gps/viewer — تعيين/قراءة الموظف المفوض."),
    (0, "GET /gps/breaches?employee_id= — المدير/المفوض — سجل الخروقات."),
    (0, "POST /gps/simulate-point — المدير/الموارد — حقن نقطة تجريبية (للاختبار على ويندوز)."),
])

add_slide("5) هاتف المدير (Manager)", [
    (0, "يعيّن الموظف المفوض من قائمة ويضغط حفظ → POST /gps/set-viewer."),
    (0, "يبدأ التتبّع: يختار الموظف + الوضع (مسافة 50–200م / زمن 1–30د) + بدء → POST /gps/start."),
    (0, "يرى الجلسات على الخريطة: دائرة المنطقة + خط السير + دبوس أحمر عند آخر نقطة."),
    (0, "شارة برتقالية 'خارج المنطقة' بعدّاد زمن/مسافة حي (يتحدث كل ثانية)."),
    (0, "أزرار: محاكاة (حقن نقطة داخل/خارج للاختبار)، إيقاف (POST /gps/stop)."),
    (0, "زر 'سجل الخروقات' → شاشة تقرير + تصدير PDF عربي."),
    (0, "يستقبل تنبيهاً فورياً (WebSocket) عند الخروج وعند العودة."),
])

add_slide("6) هاتف الموظف المُتتبَّع (Tracked Employee)", [
    (0, "لا يرى الخريطة؛ يرى شريطاً: 'يتم تتبّع موقعك' أو 'أنت خارج منطقة عملك' بعدّاد حي."),
    (0, "خدمة LocationTrackerService تعمل تلقائياً بعد الدخول (تُشغَّل في AuthProvider.login/init)."),
    (0, "تسأل إذن الموقع، وتتحقق كل 15 ثانية من GET /gps/my-active."),
    (0, "إن وُجدت جلسة نشطة: ترسل موقعها عبر POST /gps/point بالوتيرة التي حددها المدير."),
    (0, "المسافة = distanceFilter (متر)؛ الزمن = مؤقّت (دقيقة) عبر getCurrentPosition."),
    (0, "تستقبل عبر WebSocket تنبيه 'أنت خارج منطقة عملك'."),
    (0, "تتطلب أذونات: ACCESS_FINE/COARSE/BACKGROUND_LOCATION (أضيفت للـ Manifest)."),
])

add_slide("7) هاتف المفوض (Designated Viewer)", [
    (0, "موظف واحد تعيّنه الإدارة عبر AppSetting.tracking_viewer_id."),
    (0, "لا يوقف التتبّع، لكن يراقب أي موظف وطباعة مسيره (PDF)."),
    (0, "يختار موظفاً من القائمة → خريطة + عدّاد مركزي حي عند الخروج + زر 'طباعة المسير'."),
    (0, "طباعة المسير: ReportPdf.shareRoute (خط سير مرسوم + جدول إحداثيات)."),
    (0, "يستقبل تنبيهات الخروج والملخّص مثل المدير، ويفتح سجل الخروقات."),
])

add_slide("8) منطق كشف الخروج من المنطقة (Geofence)", [
    (0, "عند كل نقطة: الخادم يحسب المسافة (Haversine) بين الموقع ومركز منطقة العمل."),
    (0, "داخل المنطقة: لا إجراء."),
    (0, "خارج (أول مرة): تنبيه للموظف + المدير + المفوض؛ يبدأ عدّ الزمن والمسافة."),
    (0, "أثناء الخروج: تُجمع outside_distance_m = مجموع أطوال الأجزاء بين النقاط."),
    (0, "عند العودة: يُنشأ سجل GeofenceBreach (المدة + المسافة) + تنبيه ملخّص للمديرين والمفوض."),
    (0, "يُستثنى الموظف إن geofence_exempt=True أو لا توجد منطقة معرّفة."),
])

add_slide("9) تبادل البيانات خطوة بخطوة (Sequence)", [
    (0, "1) المدير: POST /gps/start {target, mode, interval} → الخادم ينشئ session (status=active)."),
    (0, "2) المدير: يسحب GET /gps/view كل 5ث → يظهر على الخريطة 'بانتظار أول نقطة'."),
    (0, "3) الموظف: LocationTrackerService → GET /gps/my-active → active=true + session_id + mode."),
    (0, "4) الموظف: يرسل POST /gps/point {session_id, lat, lng} كل (مسافة/زمن)."),
    (0, "5) الخادم: يخزّن FieldTrackingPoint + يشغّل كشف geofence."),
    (0, "6) خارج المنطقة → notify_employee(الموظف) + notify_managers + notify_employee(المفوض) [+ WebSocket]."),
    (0, "7) المدير/المفوض: GET /gps/view يعرض is_outside+outside_started_at+outside_distance_m (عدّاد حي)."),
    (0, "8) الموظف: GET /gps/my-active يعرض is_outside → شريط 'أنت خارج منطقة عملك' بعدّاد."),
    (0, "9) عند العودة → إنشاء GeofenceBreach + تنبيه ملخّص."),
    (0, "10) المدير: POST /gps/stop → session=completed → الموظف يتوقف عن الإرسال."),
])

# ---------- Diagram slide ----------
def add_box(s, l, t, w, h, text, fill, txtcolor=WHITE):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    shp.line.width = Pt(1.25)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = AR
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = txtcolor
    return shp


def add_conn(s, x1, y1, x2, y2, label):
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = BLUE
    conn.line.width = Pt(1.75)
    # label textbox at midpoint
    mx = (x1 + x2) / 2
    my = (y1 + y2) / 2
    tb = s.shapes.add_textbox(Inches(mx - 1.6), Inches(my - 0.3), Inches(3.2), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = AR
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = NAVY


s = prs.slides.add_slide(BLANK)
ttl = s.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.8))
tp = ttl.text_frame.paragraphs[0]
tp.alignment = PP_ALIGN.RIGHT
tr = tp.add_run()
tr.text = "10) المخطط المعماري لتدفّق البيانات"
tr.font.name = AR
tr.font.size = Pt(26)
tr.font.bold = True
tr.font.color.rgb = NAVY

# Boxes
A = add_box(s, 0.4, 1.1, 3.0, 1.3, "هاتف الموظف\n(يُرسل الموقع)", ORANGE)
B = add_box(s, 5.17, 2.9, 3.0, 1.8, "خادم FastAPI\n+ قاعدة البيانات\n(المنفذ 8001)", NAVY)
C = add_box(s, 9.9, 1.1, 3.0, 1.3, "هاتف المدير\n(يتحكّم + يراقب)", BLUE)
D = add_box(s, 9.9, 4.7, 3.0, 1.3, "هاتف المفوض\n(يراقب + يطبع)", BLUE)

# Connectors (edge midpoints)
# A right edge (3.4,1.75) -> B left edge (5.17,3.8)
add_conn(s, 3.4, 1.75, 5.17, 3.8, "POST /gps/point\n(خط العرض/الطول)")
# C left edge (9.9,1.75) -> B right edge (8.17,3.8)
add_conn(s, 9.9, 1.75, 8.17, 3.8, "POST start/stop/set-viewer")
# B top-right (8.17,2.9) -> C bottom-left (9.9,1.1+1.3=2.4)? use (9.9,2.4)
add_conn(s, 8.17, 2.9, 9.9, 2.4, "GET /gps/view + WebSocket")
# B bottom-right (8.17,4.7) -> D top-left (9.9,4.7)
add_conn(s, 8.17, 4.7, 9.9, 4.7, "WebSocket + GET /gps/view")
# B left edge (5.17,4.7) -> A bottom-left (0.4,2.4)
add_conn(s, 5.17, 4.7, 0.4, 2.4, "WebSocket تنبيه خروج")

add_slide("11) الاتصال بالشبكة (LAN مقابل إنترنت)", [
    (0, "نفس الشبكة (داخل المؤسسة): الهواتف والخادم على نفس الواي فاي؛ عنوان الخادم في ⚙ = http://IP:8001."),
    (0, "بُعد 100كم / شبكة مختلفة: الخادم يُكشَف للإنترنت."),
    (0, "للتجربة: نفق cloudflared/ngrok يكشف localhost:8001 ويعطي رابط https عاماً مؤقتاً."),
    (0, "للإنتاج: استضافة الخادم على VPS/سحابة بدومين ثابت."),
    (0, "التطبيق يحوّل http→ws و https→wss تلقائياً (لا حاجة لتعديل كود)."),
    (0, "الجدار الناري: السماح بالمنفذ 8001 دخولاً (Inbound)."),
])

add_slide("12) التنبيهات الفورية (WebSocket)", [
    (0, "المسار: /ws/notifications?token=JWT."),
    (0, "connection_manager يربط employee_id بقائمة WebSocket نشطة."),
    (0, "عند الخروج/العودة: notify_employee / notify_managers تنشئ صفّ إشعار + ترسل JSON عبر WS فوراً."),
    (0, "حمولة التنبيه: notification_id, title, message, severity, created_at."),
    (0, "التطبيق (notification_ws) يعرض الإشعار ويحدّث الشارة فوراً."),
])

add_slide("13) متطلبات التشغيل والنشر", [
    (0, "خادم يعمل دائماً (يفضّل كخدمة ويندوز لا تتوقف مع إغلاق الجلسة)."),
    (0, "فتح المنفذ 8001 بالجدار الناري (دخول)."),
    (0, "عنوان IP ثابت في الشبكة المحلية أو دومين عام عند الاستضافة."),
    (0, "إدخال منطقة عمل (قطر) لكل موظف من شاشة 'الموظفين'."),
    (0, "منح أذونات الموقع على الهواتف (السماح دائماً للخلفية)."),
    (0, "إعادة بناء APK بعد تعديل AndroidManifest (لإذن BACKGROUND_LOCATION)."),
    (0, "عنوان الخادم يُضبط من ⚙ في شاشة الدخول (افتراضي 192.168.1.175 — يُغيَّر للـ IP/الرابط الفعلي)."),
])

add_slide("14) خلاصة", [
    (0, "نظام كامل: تتبّع موقع + مسير + كشف خروق المنطقة + تنبيهات فورية + تقارير PDF."),
    (0, "المدير: يبدأ/يوقف ويعيّن المفوض.  الموظف: يُرسل موقعه تلقائياً.  المفوض: يراقب ويطبع."),
    (0, "البيانات تُخزَّن في SQLite على الخادم (محلياً) ما لم تُستضاف على الإنترنت."),
    (0, "كل التبادل عبر REST + WebSocket بصلاحيات صارمة (المدير/الموارد تبدأ؛ الموظف يرسل نقاطه فقط)."),
], note="نهاية التوثيق الشامل — نص فقط.")

prs.save(r"D:\yarmouk_water_management_pro\gps_architecture.pptx")
print("SAVED gps_architecture.pptx  slides:", len(prs.slides._sldIdLst))
